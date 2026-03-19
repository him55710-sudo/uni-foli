from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from polio_render.formats.base import BaseRenderer
from polio_render.markdown import markdown_lines_to_bullets, split_markdown_sections
from polio_render.models import RenderArtifact, RenderBuildContext
from polio_shared.paths import find_project_root


class PptxRenderer(BaseRenderer):
    extension = ".pptx"
    implementation_level = "python-pptx"

    def render(self, context: RenderBuildContext) -> RenderArtifact:
        output_path = self.prepare_output_path(context)
        self._build_presentation(context, output_path)

        relative_path = str(output_path.relative_to(find_project_root()))
        return RenderArtifact(
            absolute_path=str(output_path),
            relative_path=relative_path,
            message="PPTX presentation created with python-pptx.",
        )

    def _build_presentation(self, context: RenderBuildContext, output_path) -> None:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = context.project_title
        title_subtitle = title_slide.placeholders[1]
        title_subtitle.text = f"{context.draft_title}\nRequested by: {context.requested_by or 'anonymous'}"

        sections = split_markdown_sections(context.content_markdown)
        overview_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        overview_slide.shapes.title.text = "Draft Overview"
        overview_body = overview_slide.placeholders[1].text_frame
        overview_body.clear()

        overview_points = [title for title, _ in sections if title][:6] or ["This draft is ready for expansion."]
        for index, bullet in enumerate(overview_points):
            paragraph = overview_body.paragraphs[0] if index == 0 else overview_body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(24 if index == 0 else 20)

        for section_title, lines in sections:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = section_title or "Section"
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            bullets = markdown_lines_to_bullets(lines, max_items=7) or ["No content provided."]
            for index, bullet in enumerate(bullets):
                paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0
                paragraph.font.size = Pt(22 if index == 0 else 18)
                paragraph.font.name = "Aptos"
                paragraph.alignment = PP_ALIGN.LEFT
                if index == 0:
                    paragraph.font.bold = True

            self._decorate_slide(slide)

        presentation.save(str(output_path))

    @staticmethod
    def _decorate_slide(slide) -> None:
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(247, 248, 250)

        title = slide.shapes.title
        if title and title.text_frame and title.text_frame.paragraphs:
            first_run = title.text_frame.paragraphs[0].runs[0] if title.text_frame.paragraphs[0].runs else None
            if first_run:
                first_run.font.color.rgb = RGBColor(15, 23, 42)
                first_run.font.size = Pt(28)
                first_run.font.bold = True
